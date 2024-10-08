import os
import base64
import logging
import asyncio
import json
import requests
from typing import List, Dict, Any
from functools import lru_cache
import streamlit as st
from dotenv import load_dotenv
from gtts import gTTS
from PIL import Image
import aiohttp
import PyPDF2
import docx
import pytesseract
from deep_translator import GoogleTranslator
from datetime import datetime
from fpdf import FPDF
import aiosqlite
import re
from cachetools import TTLCache
from ratelimit import limits, sleep_and_retry
import openpyxl
from pptx import Presentation

# Load environment variables
load_dotenv()

# API Keys
GROQ_API_KEY = os.getenv("GROQ_API_KEY")
OPENAI_API_KEY = os.getenv("OPENAI_API_KEY")

# Logging configuration
logging.basicConfig(level=logging.ERROR)
logger = logging.getLogger(__name__)

# Constants
MAX_CHAT_HISTORY_LENGTH = 50
DB_PATH = "chat_history.db"
MAX_FILE_CONTENT_LENGTH = 4000
TRUNCATION_ELLIPSIS = "..."
CHUNK_SIZE = 2000

VOICE_OPTIONS = {
    "OpenAI": ["alloy", "echo", "fable", "onyx", "nova", "shimmer"],
    "gTTS": ["en", "ta", "hi"],
}

PROVIDER_OPTIONS = ["Groq", "OpenAI"]

# Import persona and content_type from the root directory
from persona import PERSONAS
from content_type import CONTENT_TYPES

# Caching for API responses
response_cache = TTLCache(maxsize=100, ttl=3600)


# Rate limiting
@sleep_and_retry
@limits(calls=5, period=60)
async def rate_limited_api_call(client, params, messages):
    return await async_stream_llm_response(client, params, messages)


# Utility Functions
def validate_prompt(prompt: str):
    if not prompt.strip():
        raise ValueError("Prompt cannot be empty")


def sanitize_input(input_text):
    return re.sub(r"[^\w\s\-.,?!]", "", input_text).strip()


def process_uploaded_file(uploaded_file):
    file_handlers = {
        "application/pdf": lambda f: " ".join(
            page.extract_text() for page in PyPDF2.PdfReader(f).pages
        ),
        "application/vnd.openxmlformats-officedocument.wordprocessingml.document": lambda f: " ".join(
            paragraph.text for paragraph in docx.Document(f).paragraphs
        ),
        "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet": lambda f: process_excel_file(
            f
        ),
        "application/vnd.ms-powerpoint": lambda f: process_ppt_file(f),
        "text/plain": lambda f: f.getvalue().decode("utf-8"),
        "text/markdown": lambda f: f.getvalue().decode("utf-8"),
        "image/jpeg": perform_ocr,
        "image/png": perform_ocr,
    }

    for file_type, handler in file_handlers.items():
        if uploaded_file.type.startswith(file_type):
            return handler(uploaded_file)

    raise ValueError("Unsupported file type")


def perform_ocr(image_file):
    try:
        image = Image.open(image_file)
        return pytesseract.image_to_string(image)
    except Exception as e:
        logger.error(f"Error performing OCR: {e}")
        raise ValueError(
            "OCR processing failed. Please ensure the image is valid and readable."
        )


def process_excel_file(file):
    wb = openpyxl.load_workbook(file)
    return "\n".join(
        " | ".join(map(str, row))
        for sheet in wb.sheetnames
        for row in wb[sheet].iter_rows(values_only=True)
    )


def process_ppt_file(file):
    prs = Presentation(file)
    return "\n".join(
        shape.text
        for slide in prs.slides
        for shape in slide.shapes
        if hasattr(shape, "text")
    )


def text_to_speech(text: str, lang: str):
    lang_map = {"English": "en", "Tamil": "ta", "Hindi": "hi"}
    lang_code = lang_map.get(lang, "en")
    tts = gTTS(text=text, lang=lang_code)
    audio_file = "temp_audio.mp3"
    tts.save(audio_file)
    with open(audio_file, "rb") as f:
        audio_bytes = f.read()
    os.remove(audio_file)
    st.session_state.audio_base64 = base64.b64encode(audio_bytes).decode()


def translate_text(text: str, target_lang: str) -> str:
    return (
        text
        if target_lang == "English"
        else GoogleTranslator(source="auto", target=target_lang).translate(text)
    )


def update_token_count(tokens: int):
    st.session_state.total_tokens = st.session_state.get("total_tokens", 0) + tokens
    st.session_state.total_cost = (
        st.session_state.get("total_cost", 0) + tokens * 0.0001
    )


def export_chat(format: str):
    chat_history = "\n\n".join(
        [
            f"**{m['role'].capitalize()}:** {m['content']}"
            for m in st.session_state.messages
        ]
    )
    timestamp = datetime.now().strftime("%Y-%m-%d_%H-%M-%S")
    filename = f"chat_exports/chat_history_{timestamp}.{format}"
    os.makedirs("chat_exports", exist_ok=True)

    if format == "md":
        with open(filename, "w") as f:
            f.write(chat_history)
        st.download_button("Download Markdown", filename, file_name=filename)
    elif format == "pdf":
        pdf = FPDF()
        pdf.add_page()
        pdf.set_font("Arial", size=12)
        pdf.multi_cell(0, 10, chat_history)
        pdf.output(filename)
        st.download_button("Download PDF", filename, file_name=filename)


async def create_database():
    async with aiosqlite.connect(DB_PATH) as db:
        await db.execute(
            """
            CREATE TABLE IF NOT EXISTS chat_history (
                id INTEGER PRIMARY KEY AUTOINCREMENT,
                timestamp DATETIME DEFAULT CURRENT_TIMESTAMP,
                chat_name TEXT,
                role TEXT,
                content TEXT
            )
        """
        )
        await db.commit()


async def save_chat_history_to_db(chat_name: str):
    async with aiosqlite.connect(DB_PATH) as db:
        await db.executemany(
            "INSERT INTO chat_history (chat_name, role, content) VALUES (?, ?, ?)",
            [
                (chat_name, message["role"], message["content"])
                for message in st.session_state.messages
            ],
        )
        await db.commit()


async def load_chat_history_from_db(chat_name: str):
    async with aiosqlite.connect(DB_PATH) as db:
        async with db.execute(
            "SELECT role, content FROM chat_history WHERE chat_name = ? ORDER BY id",
            (chat_name,),
        ) as cursor:
            st.session_state.messages = [
                {"role": row[0], "content": row[1]} async for row in cursor
            ]


async def get_saved_chat_names():
    async with aiosqlite.connect(DB_PATH) as db:
        async with db.execute("SELECT DISTINCT chat_name FROM chat_history") as cursor:
            return [row[0] async for row in cursor]


async def delete_chat(chat_name):
    async with aiosqlite.connect(DB_PATH) as db:
        await db.execute("DELETE FROM chat_history WHERE chat_name = ?", (chat_name,))
        await db.commit()


async def create_content(prompt: str, content_type: str) -> str:
    full_prompt = f"Write a {content_type} based on this prompt: {prompt}"
    generated_content = ""
    async for chunk in async_stream_llm_response(
        get_api_client(st.session_state.provider),
        st.session_state.model_params,
        [{"role": "user", "content": full_prompt}],
        st.session_state.provider,
        st.session_state.voice,
    ):
        generated_content += chunk
    return generated_content


async def summarize_text(text: str, summary_type: str) -> str:
    full_prompt = f"Please provide a {summary_type} of the following text: {text}"
    summary = ""
    async for chunk in async_stream_llm_response(
        get_api_client(st.session_state.provider),
        st.session_state.model_params,
        [{"role": "user", "content": full_prompt}],
        st.session_state.provider,
        st.session_state.voice,
    ):
        summary += chunk
    return summary


def get_model_options(provider):
    models = {
        "Groq": [
            "llama-3.1-70b-versatile",
            "llama-3.1-405b-reasoning",
            "llama-3.1-8b-instant",
            "llama3-groq-70b-8192-tool-use-preview",
            "llama3-70b-8192",
            "mixtral-8x7b-32768",
            "gemma2-9b-it",
            "whisper-large-v3",
        ],
        "OpenAI": [
            "gpt-4o-mini",
            "gpt-4o",
            "gpt-3.5-turbo",
        ],
    }
    return models.get(provider, [])


def get_max_token_limit(model):
    limits = {
        "mixtral-8x7b-32768": 32768,
        "llama-3.1-70b-versatile-131072": 131072,
        "gemma2-9b-it": 8192,
    }
    return limits.get(model, 4096)


def word_count(text):
    return len(text.split())


def initialize_session_state():
    default_values = {
        "messages": [],
        "audio_base64": "",
        "file_content": "",
        "persona": "Default",
        "model_params": {
            "model": "llama-3.1-70b-versatile",
            "max_tokens": 1024,
            "temperature": 0.5,
            "top_p": 0.9,
            "top_k": 50,
            "frequency_penalty": 0.5,
            "presence_penalty": 0.5,
        },
        "total_tokens": 0,
        "total_cost": 0,
        "enable_audio": False,
        "language": "English",
        "voice": "alloy",
        "content_creation_mode": False,
        "show_summarization": False,
        "summarization_type": "Main Takeaways",
        "content_type": "Short Story",
        "provider": "Groq" if GROQ_API_KEY else "OpenAI" if OPENAI_API_KEY else None,
        "color_scheme": "Light",
        "is_file_response_handled": False,
    }
    for key, value in default_values.items():
        st.session_state.setdefault(key, value)


def reset_current_chat():
    st.session_state.messages = []
    st.session_state.is_file_response_handled = False


def is_connected():
    try:
        return requests.get("http://www.google.com", timeout=5).status_code == 200
    except requests.ConnectionError:
        return False


def save_chat_history_locally():
    chat_data = [
        {"role": msg["role"], "content": msg["content"]}
        for msg in st.session_state.messages
    ]
    with open("local_chat_history.json", "w") as f:
        json.dump(chat_data, f)


def load_chat_history_locally():
    if os.path.exists("local_chat_history.json"):
        with open("local_chat_history.json", "r") as f:
            return json.load(f)
    return []


def save_feedback(feedback: str):
    with open("feedback.json", "a") as f:
        f.write(
            json.dumps({"feedback": feedback, "timestamp": datetime.now().isoformat()})
            + "\n"
        )


async def async_stream_groq_response(
    client, params: Dict[str, Any], messages: List[Dict[str, str]]
):
    retry_attempts = 3
    for attempt in range(retry_attempts):
        try:
            async with aiohttp.ClientSession() as session:
                for i in range(0, len(messages), CHUNK_SIZE):
                    chunk_messages = messages[i : i + CHUNK_SIZE]
                    async with session.post(
                        "https://api.groq.com/openai/v1/chat/completions",
                        headers={
                            "Authorization": f"Bearer {GROQ_API_KEY}",
                            "Content-Type": "application/json",
                        },
                        json={
                            "model": params["model"],
                            "messages": chunk_messages,
                            "max_tokens": params.get("max_tokens"),
                            "temperature": params.get("temperature"),
                            "top_p": params.get("top_p"),
                            "frequency_penalty": params.get("frequency_penalty"),
                            "presence_penalty": params.get("presence_penalty"),
                            "stream": True,
                        },
                    ) as response:
                        if response.status != 200:
                            error_text = await response.text()
                            logger.error(
                                f"Groq API Error: {response.status} - {error_text}"
                            )
                            yield f"Groq API Error: {response.status} - {error_text}"
                            return

                        async for line in response.content:
                            if line.strip() and line.startswith(b"data: "):
                                try:
                                    json_str = line[6:].decode("utf-8").strip()
                                    if json_str != "[DONE]":
                                        chunk = json.loads(json_str)
                                        if (
                                            "choices" in chunk
                                            and chunk["choices"]
                                            and "delta" in chunk["choices"][0]
                                        ):
                                            yield chunk["choices"][0]["delta"][
                                                "content"
                                            ]
                                except json.JSONDecodeError as json_err:
                                    logger.error(
                                        f"JSON decode error: {json_err}. Raw line: {line}"
                                    )
            return
        except (aiohttp.ClientError, aiohttp.ServerTimeoutError) as e:
            logger.error(f"Error in Groq API call: {str(e)}")
            if attempt < retry_attempts - 1:
                await asyncio.sleep(2)
            else:
                yield f"Error in Groq API call after retries: {str(e)}"


async def async_stream_openai_response(
    client, params: Dict[str, Any], messages: List[Dict[str, str]]
):
    try:
        if st.session_state.language != "English":
            assistant_response = translate_text(
                messages[-1]["content"], st.session_state.language
            )
            messages[-1]["content"] = assistant_response

        response = await client.chat.completions.create(
            model=params["model"],
            messages=[{"role": "assistant", "content": messages[-1]["content"]}],
            max_tokens=params["max_tokens"],
            temperature=params["temperature"],
            top_p=params["top_p"],
            frequency_penalty=params.get("frequency_penalty"),
            presence_penalty=params.get("presence_penalty"),
        )

        yield response.choices[0].message.content
    except Exception as e:
        logger.error(f"Error in OpenAI API call: {str(e)}")
        yield f"Error in OpenAI API call: {str(e)}"


async def async_stream_llm_response(
    client: Any, params: Dict[str, Any], messages: List[Dict[str, str]], provider: str
) -> str:
    try:
        if provider == "Groq":
            async for chunk in async_stream_groq_response(client, params, messages):
                yield chunk
        elif provider == "OpenAI":
            async for chunk in async_stream_openai_response(client, params, messages):
                yield chunk
        else:
            yield f"Error: Unsupported provider: {provider}"
    except Exception as e:
        logger.error(f"Unexpected error in async_stream_llm_response: {str(e)}")
        yield f"API Error: {str(e)}"


@lru_cache(maxsize=None)
def get_api_client(provider: str):
    try:
        if provider == "Groq":
            from groq import Groq

            return Groq(api_key=GROQ_API_KEY)
        elif provider == "OpenAI":
            import openai

            return openai.AsyncOpenAI(api_key=OPENAI_API_KEY)
        else:
            raise ValueError(f"Unsupported provider: {provider}")
    except Exception as e:
        logger.error(f"Error initializing {provider} client: {e}")
        return None


async def process_chat_input(prompt: str, client: Any) -> None:
    try:
        validate_prompt(prompt)

        if (
            st.session_state.file_content
            and not st.session_state.is_file_response_handled
        ):
            prompt = f"Based on the uploaded file content, {prompt}\n\nFile content: {st.session_state.file_content[:MAX_FILE_CONTENT_LENGTH]}{TRUNCATION_ELLIPSIS}"
            st.session_state.is_file_response_handled = True

        persona_content = (
            st.session_state.custom_persona
            if st.session_state.persona == "Custom"
            else PERSONAS[st.session_state.persona]
        )

        messages = [
            {"role": "system", "content": persona_content},
            *st.session_state.messages[-MAX_CHAT_HISTORY_LENGTH:],
            {"role": "user", "content": prompt},
        ]
        with st.chat_message("user"):
            st.markdown(prompt)

        full_response = ""
        with st.chat_message("assistant"):
            message_placeholder = st.empty()
            async for chunk in async_stream_llm_response(
                client,
                st.session_state.model_params,
                messages,
                st.session_state.provider,
            ):
                if chunk.startswith("API Error:"):
                    message_placeholder.error(chunk)
                    return
                full_response += chunk
                message_placeholder.markdown(full_response + "▌")

            message_placeholder.markdown(full_response)

        st.session_state.messages.extend(
            [
                {"role": "user", "content": prompt},
                {"role": "assistant", "content": full_response},
            ]
        )

        if st.session_state.enable_audio and full_response.strip():
            text_to_speech(full_response, st.session_state.language)
            st.audio(
                f"data:audio/mp3;base64,{st.session_state.audio_base64}",
                format="audio/mp3",
            )

        update_token_count(len(full_response.split()))

        if st.session_state.content_creation_mode:
            generated_content = await create_content(
                prompt, st.session_state.content_type
            )
            with st.chat_message("assistant"):
                st.markdown(
                    f"## Generated {st.session_state.content_type}:\n\n{generated_content}"
                )

        if st.session_state.show_summarization:
            summary = await summarize_text(
                st.session_state.file_content or prompt,
                st.session_state.summarization_type,
            )
            with st.chat_message("assistant"):
                st.markdown(
                    f"## Summary ({st.session_state.summarization_type}):\n\n{summary}"
                )

    except ValueError as ve:
        st.error(f"Invalid input: {str(ve)}")
    except Exception as e:
        logger.error(f"Unexpected error in process_chat_input: {str(e)}")
        st.error(f"An unexpected error occurred. Please try again later.")


def setup_sidebar() -> None:
    with st.sidebar:
        st.markdown(
            "<h2 style='text-align: center;color: #6ca395;'>Select Provider</h2>",
            unsafe_allow_html=True,
        )

        available_providers = [
            p for p in PROVIDER_OPTIONS if os.getenv(f"{p.upper()}_API_KEY")
        ]

        if not available_providers:
            st.error(
                "No API keys are set. Please set at least one API key in your .env file."
            )
            st.stop()

        selected_provider = st.selectbox(
            "Provider", available_providers, label_visibility="collapsed"
        )
        st.session_state.provider = selected_provider

        st.markdown(
            "<h2 style='text-align: center;'>Settings 🛠️ </h2>", unsafe_allow_html=True
        )

        with st.expander("Chat Settings", expanded=True):
            saved_chats = st.session_state.get("saved_chats", [])
            selected_chat = st.selectbox(
                "Load Chat History", options=[""] + saved_chats
            )
            if selected_chat:
                st.session_state.load_chat = selected_chat

            col4, col5, col6 = st.columns(3)
            with col4:
                if st.button("Reset"):
                    reset_current_chat()
            with col5:
                if st.button("New"):
                    reset_current_chat()
            with col6:
                if st.button("Delete"):
                    if selected_chat:
                        st.session_state.delete_chat = selected_chat
                        st.rerun()

            col1, col2 = st.columns([2, 1])
            with col1:
                chat_name_input = st.text_input(
                    "Enter a name for this chat:",
                    max_chars=50,
                    label_visibility="collapsed",
                    placeholder="Chat Name",
                )
            with col2:
                if st.button("Save") and chat_name_input:
                    st.session_state.save_chat = chat_name_input
                    st.rerun()

        with st.expander("Model"):
            model_options = get_model_options(st.session_state.provider)
            st.session_state.model_params["model"] = st.selectbox(
                "Choose Model:",
                options=model_options,
                index=(
                    model_options.index(st.session_state.model_params["model"])
                    if st.session_state.model_params["model"] in model_options
                    else 0
                ),
            )

            max_token_limit = get_max_token_limit(
                st.session_state.model_params["model"]
            )
            st.session_state.model_params["max_tokens"] = st.slider(
                "Max Tokens:",
                min_value=1,
                max_value=max_token_limit,
                value=min(st.session_state.model_params["max_tokens"], max_token_limit),
                step=1,
            )
            st.session_state.model_params["temperature"] = st.slider(
                "Temperature:",
                0.0,
                2.0,
                st.session_state.model_params["temperature"],
                0.1,
            )
            st.session_state.model_params["top_p"] = st.slider(
                "Top-p:", 0.0, 1.0, st.session_state.model_params["top_p"], 0.1
            )
            st.session_state.model_params["top_k"] = st.slider(
                "Top-k:", 0, 100, st.session_state.model_params["top_k"]
            )
            st.session_state.model_params["frequency_penalty"] = st.slider(
                "Frequency Penalty:",
                0.0,
                1.0,
                st.session_state.model_params["frequency_penalty"],
                0.1,
            )
            st.session_state.model_params["presence_penalty"] = st.slider(
                "Presence Penalty:",
                0.0,
                1.0,
                st.session_state.model_params["presence_penalty"],
                0.1,
            )

        with st.expander("Persona"):
            persona_options = list(PERSONAS.keys()) + ["Custom"]
            st.session_state.persona = st.selectbox(
                "Select Persona:",
                options=persona_options,
                index=persona_options.index("Default"),
            )

            if st.session_state.persona == "Custom":
                st.session_state.custom_persona = st.text_area(
                    "Enter Custom Persona Description:",
                    value=st.session_state.get("custom_persona", ""),
                    height=100,
                )
            else:
                st.text_area(
                    "Persona Description:",
                    value=PERSONAS[st.session_state.persona],
                    height=100,
                    disabled=True,
                )

        with st.expander("Audio & Language"):
            st.session_state.enable_audio = st.checkbox(
                "Enable Audio Response", value=False
            )
            st.session_state.language = st.selectbox(
                "Select Language:", ["English", "Tamil", "Hindi"]
            )
            st.session_state.voice = st.selectbox(
                "Select Voice:",
                (
                    VOICE_OPTIONS["OpenAI"]
                    if st.session_state.provider == "OpenAI"
                    else VOICE_OPTIONS["gTTS"]
                ),
            )

        with st.expander("File Upload"):
            handle_file_upload()

        with st.expander("Summarize"):
            st.session_state.show_summarization = st.checkbox(
                "Enable Summarization", value=False
            )
            if st.session_state.show_summarization:
                st.session_state.summarization_type = st.selectbox(
                    "Summarization Type:",
                    [
                        "Main Takeaways",
                        "Main points bulleted",
                        "Concise Summary",
                        "Executive Summary",
                    ],
                )

        with st.expander("Content Generation"):
            st.session_state.content_creation_mode = st.checkbox(
                "Enable Content Creation Mode", value=False
            )
            if st.session_state.content_creation_mode:
                st.session_state.content_type = st.selectbox(
                    "Select Content Type:", list(CONTENT_TYPES.keys())
                )

        with st.expander("Export"):
            export_format = st.selectbox("Export Format", ["md", "pdf"])
            st.button("Export Chat", on_click=lambda: export_chat(export_format))

        st.session_state.color_scheme = st.selectbox("Color Scheme", ["Light", "Dark"])
        if st.session_state.color_scheme == "Dark":
            st.markdown(
                "<style>.stApp { background-color: #1E1E1E; color: #FFFFFF; }</style>",
                unsafe_allow_html=True,
            )


def handle_file_upload() -> None:
    uploaded_file = st.file_uploader(
        "Upload a file",
        type=["pdf", "docx", "txt", "md", "jpg", "jpeg", "png", "xlsx", "pptx"],
    )
    if uploaded_file:
        try:
            st.session_state.file_content = process_uploaded_file(uploaded_file)
            st.success("File processed successfully")
        except Exception as e:
            st.error(f"Error processing file: {e}")


async def main() -> None:
    initialize_session_state()

    if not is_connected():
        st.warning(
            "You are currently offline. Loading previous chat history from local storage."
        )
        st.session_state.messages = load_chat_history_locally()
    else:
        await create_database()
        if "saved_chats" not in st.session_state:
            st.session_state.saved_chats = await get_saved_chat_names()

    setup_sidebar()

    if hasattr(st.session_state, "load_chat"):
        await load_chat_history_from_db(st.session_state.load_chat)
        del st.session_state.load_chat

    if hasattr(st.session_state, "save_chat"):
        await save_chat_history_to_db(st.session_state.save_chat)
        st.success(f"Chat '{st.session_state.save_chat}' saved successfully.")
        del st.session_state.save_chat

    if hasattr(st.session_state, "delete_chat"):
        await delete_chat(st.session_state.delete_chat)
        st.success(f"Chat '{st.session_state.delete_chat}' deleted successfully.")
        del st.session_state.delete_chat

    client = get_api_client(st.session_state.provider)
    if client is None:
        st.error(
            f"Failed to initialize {st.session_state.provider} client. Please check your API key."
        )
        return

    st.markdown(
        '<h1 style="text-align: center; color: #6ca395;">GenX-Chat 💬</h1>',
        unsafe_allow_html=True,
    )
    st.markdown(
        '<p style="text-align: center; color : #74a6d4">Experience the power of AI!</p>',
        unsafe_allow_html=True,
    )
    chat_container = st.container()
    with chat_container:
        for message in st.session_state.messages[-MAX_CHAT_HISTORY_LENGTH:]:
            with st.chat_message(message["role"]):
                st.markdown(message["content"])

    prompt = st.chat_input("Enter your message:")
    if prompt:
        await process_chat_input(prompt, client)

    if not is_connected() and st.session_state.messages:
        save_chat_history_locally()

    total_words = sum(word_count(msg["content"]) for msg in st.session_state.messages)
    st.sidebar.metric("Total Words", total_words)
    st.sidebar.metric("Total Tokens", st.session_state.total_tokens)
    st.sidebar.metric("Estimated Cost", f"${st.session_state.total_cost:.4f}")


if __name__ == "__main__":
    st.set_page_config(page_title="GenX-Chat", page_icon="💬", layout="wide")
    asyncio.run(main())
